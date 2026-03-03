"""XExplorer-HTML — PyQt/JS bridge.

Exposes every operation the React UI needs via pyqtSlot.
Reuses the existing database, indexer, watcher and SearchEngine backends
unchanged — only the frontend is replaced.
"""

import base64
import contextlib
import html
import json
import os
import sqlite3
import sys
import time

from PyQt6.QtCore import QFileInfo, QObject, QTimer, pyqtSignal, pyqtSlot
from PyQt6.QtGui import QPixmap
from PyQt6.QtWidgets import (
    QApplication,
    QFileDialog,
    QFileIconProvider,
    QInputDialog,
    QMessageBox,
)

from src.common.config import X_EXPLORER_DB as DB_PATH
from src.common.search_engine import SearchEngine
from src.xexplorer.database import init_db
from src.xexplorer.indexer import IndexerWorker

try:
    from watchdog.observers import Observer
    from watchdog.observers.polling import PollingObserver

    from src.xexplorer.watcher import LiveCacheUpdater
    WATCHDOG_AVAILABLE = True
except ImportError:
    WATCHDOG_AVAILABLE = False


# ── helpers ───────────────────────────────────────────────────────────────────

def _is_network_path(path: str) -> bool:
    """Return True if *path* lives on a network / UNC share."""
    if sys.platform != "win32" or not path:
        return False
    norm = path.replace("/", "\\")
    if norm.startswith("\\\\"):          # already a UNC path
        return True
    if len(norm) >= 2 and norm[1] == ":":
        drive_root = norm[0].upper() + ":\\"
        try:
            import ctypes
            DRIVE_REMOTE = 4
            return ctypes.windll.kernel32.GetDriveTypeW(drive_root) == DRIVE_REMOTE
        except Exception:
            pass
    return False


def _resolve_unc(path: str) -> str:
    if sys.platform != "win32" or not path:
        return path
    norm = path.replace("/", "\\")
    if norm.startswith("\\\\") or len(norm) < 2 or norm[1] != ":":
        return path
    drive_letter = norm[0].upper()
    drive_root = f"{drive_letter}:\\"
    try:
        import ctypes
        DRIVE_REMOTE = 4
        if ctypes.windll.kernel32.GetDriveTypeW(drive_root) != DRIVE_REMOTE:
            return path
        mpr = ctypes.WinDLL("mpr")
        buf_size = ctypes.c_ulong(1024)
        buf = ctypes.create_unicode_buffer(1024)
        ret = mpr.WNetGetUniversalNameW(f"{drive_letter}:", ctypes.c_ulong(1), buf, ctypes.byref(buf_size))
        if ret == 0:
            unc_ptr = ctypes.cast(buf, ctypes.POINTER(ctypes.c_wchar_p))[0]
            if unc_ptr:
                return unc_ptr + norm[2:]
    except Exception:
        pass
    return path


def _fmt_size(path: str, is_dir: bool) -> str:
    if is_dir:
        return ""
    try:
        b = os.path.getsize(path)
        for unit in ("B", "KB", "MB", "GB"):
            if b < 1024:
                return f"{b:.0f} {unit}"
            b /= 1024
        return f"{b:.1f} TB"
    except OSError:
        return ""


def _fmt_mtime(path: str) -> str:
    try:
        return time.strftime("%Y-%m-%d %H:%M", time.localtime(os.path.getmtime(path)))
    except OSError:
        return ""


def _db_size_mb() -> float:
    try:
        return os.path.getsize(DB_PATH) / (1024 * 1024)
    except OSError:
        return 0.0


# ── Bridge ─────────────────────────────────────────────────────────────────────────────────

class XExplorerBridge(QObject):
    # Signals → JS
    indexing_progress    = pyqtSignal(int, str)   # count, current_path
    indexing_done        = pyqtSignal(int, float) # count, seconds
    stats_updated        = pyqtSignal(str)        # JSON stats blob
    live_changed         = pyqtSignal()           # watcher detected FS change
    # Signal → Python (caught by xexplorer window to spawn child)
    open_window_requested = pyqtSignal(str)       # path to pre-navigate
    # Signal → JS (sent to a target window when a tab is dropped onto it)
    tab_incoming          = pyqtSignal(str)       # JSON {path, title}

    def __init__(self, initial_path: str = "") -> None:
        super().__init__()
        init_db()
        self._search_engine = SearchEngine(DB_PATH)
        self._worker: IndexerWorker | None = None
        self._observers: list = []
        self._icon_provider = QFileIconProvider()
        self._icon_cache: dict[str, str] = {}
        self._start_time: float = 0.0
        self._current_roots: list[str] = []
        self._initial_path: str = initial_path
        self._auto_reindex_timer: QTimer | None = None
        # Browse-folder change poller — emits live_changed when the currently
        # viewed folder's mtime changes (covers non-indexed paths too).
        self._browse_poll_path: str = ""
        self._browse_poll_mtime: float = 0.0
        self._browse_poll_timer = QTimer(self)
        self._browse_poll_timer.setInterval(750)
        self._browse_poll_timer.timeout.connect(self._poll_browse_dir)
        self._browse_poll_timer.start()

    # ── Browse-path polling ──────────────────────────────────────────────────

    @pyqtSlot(str)
    def set_active_browse_path(self, path: str) -> None:
        """JS calls this every time the active tab navigates to a new folder."""
        self._browse_poll_path = path
        try:
            self._browse_poll_mtime = os.stat(path).st_mtime if path else 0.0
        except OSError:
            self._browse_poll_mtime = 0.0

    def _poll_browse_dir(self) -> None:
        """Check if the watched browse folder changed; fire live_changed if so."""
        path = self._browse_poll_path
        if not path:
            return
        try:
            mtime = os.stat(path).st_mtime
        except OSError:
            mtime = 0.0
        if mtime != self._browse_poll_mtime:
            self._browse_poll_mtime = mtime
            self.live_changed.emit()

    # ── Window management ─────────────────────────────────────────────────────

    @pyqtSlot(result=str)
    def get_initial_path(self) -> str:
        """Return the path this window should navigate to on startup (or empty)."""
        p = self._initial_path
        self._initial_path = ""   # consume once
        return p

    @pyqtSlot(str)
    def open_new_window(self, path: str) -> None:
        """Ask the host Qt window to open a new XExplorer window at *path*."""
        self.open_window_requested.emit(path)

    @pyqtSlot()
    def close_window(self) -> None:
        """Close the host Qt window (used when the last tab is torn off)."""
        from src.xexplorer.xexplorer import _open_windows
        for win in list(_open_windows) + list(QApplication.topLevelWidgets()):
            if hasattr(win, "bridge") and win.bridge is self:
                win.close()
                return

    @pyqtSlot(str, str)
    def drop_tab(self, path: str, title: str) -> None:
        """Handle a torn-off tab drop.

        Uses QCursor.pos() (always in Qt screen coordinates, DPI-correct) to
        detect whether another XExplorer window is under the cursor.  If one is
        found it receives the tab via tab_incoming; otherwise a new window is
        spawned.
        """
        import json

        from PyQt6.QtGui import QCursor

        from src.xexplorer.xexplorer import _open_windows  # lazy to avoid circular
        cursor = QCursor.pos()
        candidates: list = []
        seen_ids: set[int] = set()
        for win in list(_open_windows) + list(QApplication.topLevelWidgets()):
            if not hasattr(win, "bridge"):
                continue
            if id(win) in seen_ids:
                continue
            seen_ids.add(id(win))
            if win.bridge is self:
                continue  # skip the source window
            if not win.isVisible():
                continue
            candidates.append(win)
            if win.frameGeometry().contains(cursor) or win.geometry().contains(win.mapFromGlobal(cursor)):
                win.bridge.tab_incoming.emit(json.dumps({"path": path, "title": title}))
                return
        # Fallback: snap to the nearest other window if cursor is very close
        # (helps with multi-monitor / DPI coordinate edge cases).
        nearest = None
        nearest_dist = 10**9
        for win in candidates:
            rect = win.frameGeometry()
            dx = max(rect.left() - cursor.x(), 0, cursor.x() - rect.right())
            dy = max(rect.top() - cursor.y(), 0, cursor.y() - rect.bottom())
            dist = dx + dy
            if dist < nearest_dist:
                nearest_dist = dist
                nearest = win
        if nearest is not None and nearest_dist <= 80:
            nearest.bridge.tab_incoming.emit(json.dumps({"path": path, "title": title}))
            return
        # No target window found — open a brand-new window
        self.open_window_requested.emit(path)

    # ── Config ───────────────────────────────────────────────────────────────

    @pyqtSlot(result=str)
    def get_config(self) -> str:
        """Return {folders:[{path,label}], ignore:[{rule,enabled}]}"""
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            # Folders
            c.execute("SELECT value FROM settings WHERE key='folders'")
            res = c.fetchone()
            folders: list[dict] = []
            if res:
                try:
                    for f in json.loads(res[0]):
                        path  = _resolve_unc(f.get("path", ""))
                        label = f.get("label", path)
                        folders.append({"path": path, "label": label})
                except (json.JSONDecodeError, TypeError):
                    pass

            # Ignore rules
            win_dir    = os.environ.get("SYSTEMROOT",       "C:\\Windows")
            prog_files = os.environ.get("PROGRAMFILES",     "C:\\Program Files")
            prog_x86   = os.environ.get("PROGRAMFILES(X86)","C:\\Program Files (x86)")
            defaults = [
                "node_modules","venv",".venv","env","__pycache__",".git",".svn",
                ".idea",".vscode","dist","build","AppData","Local Settings",
                "System Volume Information","$RECYCLE.BIN",
                ".exe",".dll",".sys",".tmp",".pyc",
                win_dir, prog_files, prog_x86,
                "C:\\MSOCache","C:\\$Recycle.Bin",
            ]
            c.execute("SELECT value FROM settings WHERE key='ignore'")
            res2 = c.fetchone()
            current: dict[str, bool] = {}
            if res2:
                for raw in res2[0].split("|"):
                    if ":" in raw:
                        rule, st = raw.rsplit(":", 1)
                        current[rule] = st == "1"
                    elif raw:
                        current[raw] = True
            for d in defaults:
                if d not in current:
                    current[d] = True
            ignore = [{"rule": k, "enabled": v} for k, v in sorted(current.items(), key=lambda x: x[0].lower())]

        # Start filesystem watchers on first config load (and whenever config is re-read)
        if folders and not self._observers:
            self._restart_watchers(folders)
            self._schedule_auto_reindex(folders)

        return json.dumps({"folders": folders, "ignore": ignore})

    @pyqtSlot(str)
    def save_config(self, json_str: str) -> None:
        data = json.loads(json_str)
        folders_json = json.dumps(data.get("folders", []))
        ignore_parts  = [f"{r['rule']}:{'1' if r['enabled'] else '0'}" for r in data.get("ignore", [])]
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute("INSERT OR REPLACE INTO settings VALUES(?,?)", ("folders", folders_json))
            c.execute("INSERT OR REPLACE INTO settings VALUES(?,?)", ("ignore",  "|".join(ignore_parts)))
            conn.commit()
        self._restart_watchers(data.get("folders", []))

    # ── Folder / drive management ─────────────────────────────────────────────

    @pyqtSlot(result=str)
    def pick_folder(self) -> str:
        path = QFileDialog.getExistingDirectory(None, "Select Folder to Index")
        if path:
            path = path.replace("/", "\\")
            return _resolve_unc(path)
        return ""

    @pyqtSlot(result=str)
    def get_drives(self) -> str:
        drives = []
        if sys.platform == "win32":
            try:
                import ctypes
                bitmask = ctypes.windll.kernel32.GetLogicalDrives()
                for letter in "ABCDEFGHIJKLMNOPQRSTUVWXYZ":
                    if bitmask & (1 << (ord(letter) - 65)):
                        d = f"{letter}:\\"
                        drives.append(_resolve_unc(d))
            except Exception:
                drives = []
        else:
            drives = ["/"]
        return json.dumps(drives)

    @pyqtSlot(str, result=str)
    def list_folder(self, path: str) -> str:
        """Return JSON list of items in a directory (same schema as search)."""
        if not path or not os.path.isdir(path):
            return json.dumps([])
        items: list[dict] = []
        try:
            with os.scandir(path) as it:
                entries = sorted(it, key=lambda e: (not e.is_dir(follow_symlinks=False), e.name.lower()))
            for entry in entries:
                try:
                    is_dir = entry.is_dir(follow_symlinks=False)
                    ext = "" if is_dir else os.path.splitext(entry.name)[1].lstrip(".").lower()
                    items.append({
                        "path":   entry.path,
                        "name":   entry.name,
                        "is_dir": is_dir,
                        "ext":    ext,
                        "size":   _fmt_size(entry.path, is_dir),
                        "mtime":  _fmt_mtime(entry.path),
                    })
                except (OSError, PermissionError):
                    continue
        except (OSError, PermissionError):
            return json.dumps([])
        return json.dumps(items)

    @pyqtSlot(str, result=str)
    def get_drive_info(self, path: str) -> str:
        """Return JSON {label, total_gb, free_gb, used_pct} for a path."""
        info: dict = {"label": "", "total_gb": 0.0, "free_gb": 0.0, "used_pct": 0}
        if sys.platform == "win32":
            try:
                import ctypes
                free_bytes  = ctypes.c_ulonglong(0)
                total_bytes = ctypes.c_ulonglong(0)
                ctypes.windll.kernel32.GetDiskFreeSpaceExW(
                    ctypes.c_wchar_p(path), None,
                    ctypes.byref(total_bytes),
                    ctypes.byref(free_bytes),
                )
                total_gb = total_bytes.value / (1024 ** 3)
                free_gb  = free_bytes.value  / (1024 ** 3)
                used_pct = round((1 - free_gb / total_gb) * 100) if total_gb else 0
                label_buf = ctypes.create_unicode_buffer(256)
                ctypes.windll.kernel32.GetVolumeInformationW(
                    ctypes.c_wchar_p(path), label_buf, 256,
                    None, None, None, None, 0,
                )
                info = {
                    "label":    label_buf.value,
                    "total_gb": round(total_gb, 1),
                    "free_gb":  round(free_gb,  1),
                    "used_pct": used_pct,
                }
            except Exception:
                pass
        return json.dumps(info)

    @pyqtSlot(str)
    def add_ignore_rule(self, rule: str) -> None:
        """Append a new ignore rule to the persisted config (state=enabled)."""
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute("SELECT value FROM settings WHERE key='ignore'")
            res = c.fetchone()
            current = res[0] if res else ""
            # Avoid duplicates
            existing = {r.rsplit(":", 1)[0] for r in current.split("|") if r}
            if rule not in existing:
                new_val = (current + "|" if current else "") + f"{rule}:1"
                c.execute("INSERT OR REPLACE INTO settings VALUES(?,?)", ("ignore", new_val))
                conn.commit()

    @pyqtSlot(result=str)
    def prompt_ignore_rule(self) -> str:
        """Open a Qt dialog to ask for a rule string; returns it or empty string."""
        rule, ok = QInputDialog.getText(None, "Add Ignore Rule", "Folder name, extension, or path:")
        return rule if ok and rule else ""

    # ── Favorites ─────────────────────────────────────────────────────────────

    @pyqtSlot(result=str)
    def get_favorites(self) -> str:
        """Return JSON list of [{path, label, icon}] favourites."""
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute("SELECT value FROM settings WHERE key='favorites'")
            res = c.fetchone()
            if res:
                try:
                    return res[0]
                except Exception:
                    pass
        return "[]"

    @pyqtSlot(str)
    def save_favorites(self, json_str: str) -> None:
        """Persist the favorites list."""
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute("INSERT OR REPLACE INTO settings VALUES(?,?)", ("favorites", json_str))
            conn.commit()

    # ── Search ────────────────────────────────────────────────────────────────

    @pyqtSlot(str, str, str, result=str)
    def search(self, query: str, filter_type: str, folders_json: str) -> str:
        """Return JSON list of {path, name, is_dir, size, mtime, ext}."""
        folder_paths: list[str] = json.loads(folders_json) if folders_json else []
        if not folder_paths:
            return "[]"
        terms = query.strip().split()
        if len(query.strip()) < 2:
            return "[]"

        if filter_type == "content":
            raw = self._search_engine.search_content(query_terms=terms, target_folders=folder_paths)
            pairs = [(r[0], r[1]) for r in raw]
        else:
            raw = self._search_engine.search_files(
                query_terms=terms,
                target_folders=folder_paths,
                files_only=(filter_type == "files"),
                folders_only=(filter_type == "folders"),
            )
            pairs = [(r[0], r[1]) for r in raw]

        results = []
        for path, is_dir in pairs[:3000]:
            if not os.path.exists(path):
                continue
            name = os.path.basename(path) or path
            _, ext = os.path.splitext(name)
            results.append({
                "path":   path,
                "name":   name,
                "is_dir": bool(is_dir),
                "size":   _fmt_size(path, is_dir),
                "mtime":  _fmt_mtime(path),
                "ext":    ext.lower().lstrip("."),
            })
        return json.dumps(results)

    # ── Indexing ──────────────────────────────────────────────────────────────

    @pyqtSlot(str)
    def start_indexing(self, paths_json: str) -> None:
        roots: list[str] = json.loads(paths_json)
        if not roots:
            return
        ignore_rules = self._get_active_ignore_rules()
        self._current_roots = roots
        self._start_time = time.time()
        self._worker = IndexerWorker(roots, ignore_rules)
        self._worker.progress.connect(self._on_progress)
        self._worker.finished.connect(self._on_done)
        self._worker.start()

    @pyqtSlot()
    def stop_indexing(self) -> None:
        if self._worker and self._worker.isRunning():
            self._worker.stop()

    def _get_active_ignore_rules(self) -> list[str]:
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            c.execute("SELECT value FROM settings WHERE key='ignore'")
            res = c.fetchone()
        rules = []
        if res:
            for raw in res[0].split("|"):
                if ":" in raw:
                    rule, st = raw.rsplit(":", 1)
                    if st == "1":
                        rules.append(rule)
                elif raw:
                    rules.append(raw)
        return rules

    def _on_progress(self, count: int, msg: str) -> None:
        self.indexing_progress.emit(count, msg)

    def _on_done(self, count: int) -> None:
        dur = time.time() - self._start_time
        now_str = time.strftime("%Y-%m-%d %H:%M:%S")
        with sqlite3.connect(DB_PATH) as conn:
            c = conn.cursor()
            for root in self._current_roots:
                c.execute("INSERT OR REPLACE INTO folder_stats VALUES(?,?)", (root, now_str))
            c.execute("INSERT OR REPLACE INTO settings VALUES(?,?)", ("last_indexed", now_str))
            conn.commit()
        self.indexing_done.emit(count, dur)
        self._emit_stats()
        self._schedule_auto_reindex()  # start periodic timer after first index completes

    # ── Periodic auto-reindex ────────────────────────────────────────────────

    _AUTO_REINDEX_INTERVAL_MS = 30 * 60 * 1000  # 30 minutes

    def _schedule_auto_reindex(self, folders: list[dict] | None = None) -> None:
        """Start a repeating timer that re-indexes all configured folders periodically.
        Safe to call multiple times — only one timer is ever active.
        """
        if self._auto_reindex_timer is not None:
            return  # already scheduled
        self._auto_reindex_timer = QTimer(self)
        self._auto_reindex_timer.setInterval(self._AUTO_REINDEX_INTERVAL_MS)
        self._auto_reindex_timer.timeout.connect(self._auto_reindex)
        self._auto_reindex_timer.start()

    def _auto_reindex(self) -> None:
        """Silently re-index all configured folders (skipped if a manual index is running)."""
        if self._worker and self._worker.isRunning():
            return  # don't interrupt a user-triggered index
        # Reload folder list from DB so we always use the current config
        try:
            with sqlite3.connect(DB_PATH) as conn:
                c = conn.cursor()
                c.execute("SELECT value FROM settings WHERE key='folders'")
                res = c.fetchone()
            roots = [f["path"] for f in json.loads(res[0])] if res else []
        except Exception:
            roots = []
        if not roots:
            return
        roots = [r for r in roots if os.path.isdir(r)]
        if not roots:
            return
        ignore_rules = self._get_active_ignore_rules()
        self._current_roots = roots
        self._start_time = time.time()
        self._worker = IndexerWorker(roots, ignore_rules)
        self._worker.progress.connect(self._on_progress)
        self._worker.finished.connect(self._on_done)
        self._worker.start()

    # ── Stats ─────────────────────────────────────────────────────────────────

    @pyqtSlot(result=str)
    def get_stats(self) -> str:
        return self._build_stats_json()

    def _build_stats_json(self) -> str:
        try:
            with sqlite3.connect(DB_PATH) as conn:
                c = conn.cursor()
                c.execute("SELECT COUNT(*) FROM files")
                count = c.fetchone()[0]
                c.execute("SELECT value FROM settings WHERE key='last_indexed'")
                res = c.fetchone()
                last = res[0] if res else "Never"
            return json.dumps({"count": count, "last_indexed": last, "db_mb": round(_db_size_mb(), 1)})
        except Exception:
            return json.dumps({"count": 0, "last_indexed": "Never", "db_mb": 0})

    def _emit_stats(self) -> None:
        self.stats_updated.emit(self._build_stats_json())

    # ── Clear DB ──────────────────────────────────────────────────────────────

    @pyqtSlot()
    def clear_index(self) -> None:
        with sqlite3.connect(DB_PATH) as conn:
            conn.execute("DELETE FROM files")
            conn.execute("DELETE FROM folder_stats")
            conn.commit()
            conn.execute("VACUUM")
            conn.commit()
        self._emit_stats()

    # ── File operations ───────────────────────────────────────────────────────

    @pyqtSlot(str)
    def open_path(self, path: str) -> None:
        if os.path.exists(path):
            os.startfile(path)

    @pyqtSlot(str)
    def show_in_explorer(self, path: str) -> None:
        d = path if os.path.isdir(path) else os.path.dirname(path)
        if os.path.exists(d):
            if sys.platform == "win32":
                import subprocess
                subprocess.Popen(["explorer", "/select,", path] if os.path.isfile(path) else ["explorer", d])
            else:
                os.startfile(d)

    @pyqtSlot(str)
    def copy_to_clipboard(self, text: str) -> None:
        QApplication.clipboard().setText(text)

    @pyqtSlot(str)
    def open_in_file_ops(self, paths_json: str) -> None:
        """Open File Tools window with the given paths pre-loaded."""
        try:
            from src.file_ops.file_ops import FileToolsWindow
            paths = json.loads(paths_json)
            win = FileToolsWindow()
            win.source_paths = list(paths)
            win._refresh_list()
            win.show()
            self.__file_ops_win = win  # keep alive
        except Exception as e:
            QMessageBox.warning(None, "File Ops Error", str(e))

    @pyqtSlot(str)
    def open_in_archiver(self, paths_json: str) -> None:
        """Open Archiver window with the given paths pre-loaded."""
        try:
            from src.archiver.archiver import ArchiverWindow
            paths = json.loads(paths_json)
            win = ArchiverWindow()
            win.source_paths = list(paths)
            win._refresh_list()
            win.show()
            self.__archiver_win = win  # keep alive
        except Exception as e:
            QMessageBox.warning(None, "Archiver Error", str(e))

    @pyqtSlot(str, result=str)
    def get_file_icon_b64(self, path: str) -> str:
        """Return a base64-encoded PNG of the system file icon (~24px)."""
        ext = os.path.splitext(path)[1].lower() if not os.path.isdir(path) else "__DIR__"
        if ext in self._icon_cache:
            return self._icon_cache[ext]
        try:
            fi = QFileInfo(path) if os.path.exists(path) else None
            icon = self._icon_provider.icon(fi) if fi else (
                self._icon_provider.icon(QFileIconProvider.IconType.Folder)
                if os.path.isdir(path)
                else self._icon_provider.icon(QFileIconProvider.IconType.File)
            )
            pm: QPixmap = icon.pixmap(24, 24)
            buf = pm.toImage()
            # Write via QBuffer workaround
            from PyQt6.QtCore import QBuffer, QIODevice
            ba = bytearray()
            buf2 = QBuffer()
            buf2.open(QIODevice.OpenMode.WriteOnly)
            buf.save(buf2, "PNG")
            data = bytes(buf2.data())
            b64 = base64.b64encode(data).decode()
            self._icon_cache[ext] = b64
            return b64
        except Exception:
            return ""

    @pyqtSlot(str, result=str)
    def get_preview(self, path: str) -> str:
        """Return JSON {type, content, ...} for the preview pane (first page/slide/sheet)."""
        if not os.path.exists(path) or os.path.isdir(path):
            return json.dumps({"type": "unsupported", "content": ""})

        ext = os.path.splitext(path)[1].lower()
        IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".svg", ".ico"}
        TEXT_EXTS  = {
            ".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx", ".json", ".yaml",
            ".yml", ".toml", ".ini", ".cfg", ".html", ".css", ".xml", ".csv",
            ".sh", ".bat", ".ps1", ".rs", ".go", ".java", ".c", ".cpp", ".h",
            ".log", ".env", ".gitignore",
        }
        PDF_EXTS    = {".pdf"}
        OFFICE_EXTS = {".pptx", ".ppt", ".xlsx", ".xls", ".xlsm", ".docx"}

        # Delegate paginated types to get_preview_page(path, 0)
        if ext in PDF_EXTS or ext in OFFICE_EXTS:
            return self.get_preview_page(path, 0)

        try:
            size = os.path.getsize(path)
            if ext in IMAGE_EXTS:
                if size > 8 * 1024 * 1024:  # 8 MB cap
                    return json.dumps({"type": "unsupported", "content": "Image too large to preview."})
                with open(path, "rb") as f:
                    data = base64.b64encode(f.read()).decode()
                mime = {".svg": "image/svg+xml", ".gif": "image/gif"}.get(ext, "image/png")
                return json.dumps({"type": "image", "content": f"data:{mime};base64,{data}", "ext": ext})
            if ext in TEXT_EXTS or size < 256 * 1024:
                with open(path, encoding="utf-8", errors="replace") as f:
                    content = f.read(64 * 1024)
                return json.dumps({"type": "text", "content": content, "ext": ext.lstrip(".")})
        except Exception as e:
            return json.dumps({"type": "error", "content": str(e)})
        return json.dumps({"type": "unsupported", "content": ""})

    @pyqtSlot(str, int, result=str)
    def get_preview_page(self, path: str, page: int) -> str:
        """Return preview for a specific page/slide/sheet (0-based).

        Handles: PDF (.pdf), PowerPoint (.pptx/.ppt), Excel (.xlsx/.xls/.xlsm), Word (.docx).
        """
        if not os.path.exists(path) or os.path.isdir(path):
            return json.dumps({"type": "unsupported", "content": ""})

        ext = os.path.splitext(path)[1].lower()
        try:
            # ── PDF ──────────────────────────────────────────────────────────
            if ext == ".pdf":
                try:
                    import fitz  # PyMuPDF
                    doc = fitz.open(path)
                    n   = len(doc)
                    idx = max(0, min(page, n - 1))
                    mat = fitz.Matrix(1.5, 1.5)
                    pix = doc[idx].get_pixmap(matrix=mat)
                    data = base64.b64encode(pix.tobytes("png")).decode()
                    doc.close()
                    return json.dumps({
                        "type": "pdf",
                        "content": f"data:image/png;base64,{data}",
                        "page_count": n,
                        "page": idx,
                        "label": f"Page {idx + 1} of {n}",
                    })
                except ImportError:
                    return json.dumps({"type": "unsupported", "content": "Install PyMuPDF:\npip install pymupdf"})

            # ── PowerPoint ───────────────────────────────────────────────────
            if ext in {".pptx", ".ppt"}:
                # Preferred path on Windows: export real slide images via
                # PowerPoint COM so the user sees actual slides (layout/colors).
                if sys.platform == "win32":
                    try:
                        import hashlib
                        import tempfile

                        import pythoncom
                        import win32com.client

                        pythoncom.CoInitialize()
                        app = None
                        prs = None
                        try:
                            app = win32com.client.DispatchEx("PowerPoint.Application")
                            prs = app.Presentations.Open(path, WithWindow=False, ReadOnly=True)
                            n = int(prs.Slides.Count)
                            idx = max(0, min(page, n - 1))

                            st = os.stat(path)
                            key_src = f"{os.path.abspath(path)}|{st.st_size}|{st.st_mtime_ns}"
                            key = hashlib.sha1(key_src.encode("utf-8", errors="replace")).hexdigest()
                            cache_dir = os.path.join(tempfile.gettempdir(), "xexplorer_ppt_cache", key)
                            os.makedirs(cache_dir, exist_ok=True)

                            img_path = os.path.join(cache_dir, f"slide_{idx + 1:04d}.png")
                            if not os.path.exists(img_path):
                                prs.Slides(idx + 1).Export(img_path, "PNG", 1600, 900)

                            with open(img_path, "rb") as f:
                                data = base64.b64encode(f.read()).decode()

                            title_text = ""
                            with contextlib.suppress(Exception):
                                title_shape = prs.Slides(idx + 1).Shapes.Title
                                if title_shape and title_shape.TextFrame and title_shape.TextFrame.HasText:
                                    title_text = str(title_shape.TextFrame.TextRange.Text or "").strip()

                            return json.dumps({
                                "type": "slide_image",
                                "content": f"data:image/png;base64,{data}",
                                "page_count": n,
                                "page": idx,
                                "label": f"Slide {idx + 1}/{n}: {title_text}" if title_text else f"Slide {idx + 1}/{n}",
                            })
                        finally:
                            with contextlib.suppress(Exception):
                                if prs is not None:
                                    prs.Close()
                            with contextlib.suppress(Exception):
                                if app is not None:
                                    app.Quit()
                            with contextlib.suppress(Exception):
                                pythoncom.CoUninitialize()
                    except Exception:
                        # Fall through to text-based fallback below.
                        pass

                try:
                    from pptx import Presentation  # python-pptx
                    prs    = Presentation(path)
                    slides = prs.slides
                    n      = len(slides)
                    idx    = max(0, min(page, n - 1))
                    slide  = slides[idx]

                    def _slide_html(s) -> str:
                        parts: list[str] = []
                        for shape in s.shapes:
                            if not shape.has_text_frame:
                                continue
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if not t:
                                    continue
                                lvl   = para.level
                                tag   = "h2" if lvl == 0 and not parts else "p"
                                style = f"margin-left:{lvl * 16}px"
                                parts.append(f'<{tag} class="sl-text sl-l{lvl}" style="{style}">{t}</{tag}>')
                        return "\n".join(parts) or '<p class="sl-empty">Empty slide</p>'

                    title_text = ""
                    if slide.shapes.title:
                        title_text = slide.shapes.title.text or ""
                    return json.dumps({
                        "type":       "slide",
                        "content":    _slide_html(slide),
                        "page_count": n,
                        "page":       idx,
                        "label":      f"Slide {idx + 1}/{n}: {title_text}" if title_text else f"Slide {idx + 1}/{n}",
                    })
                except ImportError:
                    return json.dumps({"type": "unsupported", "content": "Install python-pptx:\npip install python-pptx"})

            # ── Excel ────────────────────────────────────────────────────────
            if ext in {".xlsx", ".xls", ".xlsm"}:
                try:
                    import openpyxl
                    wb     = openpyxl.load_workbook(path, read_only=True, data_only=True)
                    sheets = wb.sheetnames
                    n      = len(sheets)
                    idx    = max(0, min(page, n - 1))
                    ws     = wb[sheets[idx]]
                    rows: list[tuple] = []
                    for i, row in enumerate(ws.iter_rows(values_only=True)):
                        if i >= 200:
                            break
                        rows.append(row)
                    wb.close()

                    def _table_html(rows: list[tuple]) -> str:
                        if not rows:
                            return '<p class="sl-empty">Empty sheet</p>'
                        head = rows[0]
                        html = '<div class="sheet-wrap"><table class="sheet-table"><thead><tr>'
                        for cell in head:
                            v = "" if cell is None else str(cell)
                            html += f"<th>{v}</th>"
                        html += "</tr></thead><tbody>"
                        for row in rows[1:]:
                            html += "<tr>"
                            for cell in row:
                                v = "" if cell is None else str(cell)
                                html += f"<td>{v}</td>"
                            html += "</tr>"
                        html += "</tbody></table></div>"
                        return html

                    return json.dumps({
                        "type":       "sheet",
                        "content":    _table_html(rows),
                        "page_count": n,
                        "page":       idx,
                        "label":      f"{sheets[idx]} ({idx + 1}/{n})",
                    })
                except ImportError:
                    return json.dumps({"type": "unsupported", "content": "Install openpyxl:\npip install openpyxl"})

            # ── Word (DOCX) ───────────────────────────────────────────────────
            if ext == ".docx":
                try:
                    from docx import Document
                    doc = Document(path)
                    blocks: list[str] = []

                    for p in doc.paragraphs:
                        text = (p.text or "").strip()
                        if not text:
                            continue
                        style_name = (p.style.name or "").lower() if p.style else ""
                        tag = "p"
                        if style_name.startswith("heading 1"):
                            tag = "h1"
                        elif style_name.startswith("heading 2"):
                            tag = "h2"
                        elif style_name.startswith("heading 3"):
                            tag = "h3"
                        elif "title" in style_name:
                            tag = "h1"
                        elif "subtitle" in style_name:
                            tag = "h2"
                        safe = html.escape(text)
                        blocks.append(f"<{tag}>{safe}</{tag}>")

                    for table in doc.tables:
                        rows_html: list[str] = []
                        for row in table.rows:
                            cells = [html.escape((c.text or "").strip()) for c in row.cells]
                            row_html = "".join(f"<td>{c}</td>" for c in cells)
                            rows_html.append(f"<tr>{row_html}</tr>")
                        if rows_html:
                            blocks.append(
                                '<div class="sheet-wrap"><table class="sheet-table"><tbody>'
                                + "".join(rows_html)
                                + "</tbody></table></div>"
                            )

                    if not blocks:
                        blocks.append('<p class="sl-empty">Empty document</p>')

                    return json.dumps({
                        "type": "docx",
                        "content": "\n".join(blocks),
                        "page_count": 1,
                        "page": 0,
                        "label": "Word document",
                    })
                except ImportError:
                    return json.dumps({"type": "unsupported", "content": "Install python-docx:\npip install python-docx"})

        except Exception as e:
            return json.dumps({"type": "error", "content": str(e)})

        return json.dumps({"type": "unsupported", "content": ""})

    # ── Live watcher ──────────────────────────────────────────────────────────

    def _restart_watchers(self, folders: list[dict]) -> None:
        if not WATCHDOG_AVAILABLE:
            return
        # Stop any previously running observers
        for obs in self._observers:
            with contextlib.suppress(Exception):
                obs.stop()
                obs.join()
        self._observers.clear()

        paths = [f["path"] for f in folders if os.path.exists(f.get("path", ""))]
        if not paths:
            return

        ignore_rules = self._get_active_ignore_rules()
        handler = LiveCacheUpdater(ignore_rules)
        handler.file_changed = lambda: self.live_changed.emit()  # type: ignore

        # Split into local vs. network paths so each gets the right observer
        local_paths   = [p for p in paths if not _is_network_path(p)]
        network_paths = [p for p in paths if     _is_network_path(p)]

        def _make_observer(path_list, observer_cls, label):
            if not path_list:
                return
            obs = observer_cls()
            scheduled = 0
            for p in path_list:
                with contextlib.suppress(Exception):
                    obs.schedule(handler, p, recursive=True)
                    scheduled += 1
            if scheduled:
                obs.start()
                self._observers.append(obs)

        _make_observer(local_paths,   Observer,        "native")
        # Poll every 10 s for network paths — fast enough without hammering the share
        _make_observer(network_paths, lambda: PollingObserver(timeout=10), "polling")


    @pyqtSlot(result=bool)
    def is_watchdog_available(self) -> bool:
        return WATCHDOG_AVAILABLE
